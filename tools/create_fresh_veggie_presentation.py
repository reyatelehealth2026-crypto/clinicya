from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(r"C:\Users\Administrator\clinicya\REYA_Fresh_Veggie_Model_Gap_Analysis.pptx")

W, H = Inches(13.333), Inches(7.5)

GREEN = RGBColor(46, 125, 50)
DARK_GREEN = RGBColor(27, 94, 32)
SOFT_GREEN = RGBColor(232, 245, 233)
MINT = RGBColor(224, 247, 239)
ORANGE = RGBColor(245, 124, 0)
SOFT_ORANGE = RGBColor(255, 243, 224)
BLUE = RGBColor(25, 103, 210)
SOFT_BLUE = RGBColor(232, 240, 254)
PURPLE = RGBColor(93, 63, 211)
SOFT_PURPLE = RGBColor(241, 238, 255)
RED = RGBColor(198, 40, 40)
SOFT_RED = RGBColor(255, 235, 238)
GRAY = RGBColor(86, 86, 86)
LIGHT_GRAY = RGBColor(248, 249, 250)
MID_GRAY = RGBColor(220, 224, 228)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(24, 24, 24)

FONT = "Leelawadee UI"
FONT_BOLD = "Leelawadee UI"


def set_text(tf, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BOLD if bold else FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def add_textbox(slide, x, y, w, h, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.margin_left = Inches(0.03)
    box.text_frame.margin_right = Inches(0.03)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.word_wrap = True
    set_text(box.text_frame, text, size, color, bold, align)
    return box


def add_card(slide, x, y, w, h, title, body="", fill=WHITE, line=MID_GRAY, title_color=BLACK,
             icon="", radius=True, title_size=16, body_size=11):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = (icon + "  " if icon else "") + title
    r.font.name = FONT_BOLD
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = title_color
    if body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = FONT
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = GRAY
    return shp


def add_header(slide, title, subtitle=None):
    add_textbox(slide, 0.45, 0.28, 8.9, 0.42, title, 23, DARK_GREEN, True)
    if subtitle:
        add_textbox(slide, 0.47, 0.74, 9.4, 0.32, subtitle, 11.5, GRAY)
    add_textbox(slide, 10.3, 0.32, 2.5, 0.25, "REYA / Fresh Veggie Model", 9, GREEN, True, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(1.05), Inches(12.45), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(204, 232, 207)
    line.line.fill.background()


def add_footer(slide, page):
    add_textbox(slide, 0.45, 7.12, 5.2, 0.22, "Business model gap analysis: ธุรกิจผักสดผ่าน LINE + Mini App", 8.5, GRAY)
    add_textbox(slide, 12.15, 7.12, 0.75, 0.22, str(page), 8.5, GRAY, False, PP_ALIGN.RIGHT)


def add_leaf_decor(slide):
    for x, y, s, c in [(12.15, 0.14, 0.33, SOFT_GREEN), (12.55, 0.42, 0.22, SOFT_ORANGE), (0.18, 6.84, 0.25, SOFT_GREEN)]:
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s * 0.55))
        shp.rotation = -28
        shp.fill.solid()
        shp.fill.fore_color.rgb = c
        shp.line.color.rgb = RGBColor(190, 220, 190)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.4, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    if dash:
        conn.line.dash_style = 4
    return conn


def bullet_box(slide, x, y, w, title, items, color=GREEN, fill=SOFT_GREEN, icon=""):
    add_card(slide, x, y, w, 0.55, title, fill=fill, line=color, title_color=color, icon=icon, title_size=15)
    top = y + 0.68
    for i, text in enumerate(items):
        add_textbox(slide, x + 0.08, top + i * 0.38, w - 0.16, 0.28, "• " + text, 10.5, BLACK)


def status_pill(slide, x, y, text, fill, line, color=BLACK):
    return add_card(slide, x, y, 1.12, 0.28, text, fill=fill, line=line, title_color=color, title_size=8.5)


def create_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 253, 248)
    add_leaf_decor(slide)
    add_textbox(slide, 0.65, 0.55, 4.2, 0.55, "Fresh Veggie Delivery", 27, ORANGE, True)
    add_textbox(slide, 0.68, 1.12, 6.1, 0.58, "แผนปรับโปรเจค REYA/clinicya เป็นธุรกิจผักสด", 25, DARK_GREEN, True)
    add_textbox(slide, 0.70, 1.82, 5.7, 0.42, "วิเคราะห์ระบบที่มีแล้ว / ระบบที่ยังขาด / Roadmap เพื่อทำโมเดลแบบ LINE + Mini App + Subscription + Delivery", 13, GRAY)
    add_card(slide, 0.72, 2.55, 3.6, 1.1, "มีฐานพร้อม", "LINE OA, Inbox, Mini App, Customer, Order, Admin", SOFT_GREEN, GREEN, DARK_GREEN, "✓")
    add_card(slide, 4.65, 2.55, 3.6, 1.1, "จุดที่ยังขาด", "Fresh stock, Subscription, Delivery, Proof, Claim", SOFT_ORANGE, ORANGE, ORANGE, "!")
    add_card(slide, 8.58, 2.55, 3.6, 1.1, "เป้าหมาย", "ขายผักสดแบบสั่งครั้งเดียว + กล่องผักประจำ + ติดตามสถานะ", SOFT_BLUE, BLUE, BLUE, "→")
    add_card(slide, 0.72, 4.15, 11.55, 1.3, "Executive take-away", "ระบบเดิมเหมาะกับร้านค้าออนไลน์และ CRM แล้ว แต่ธุรกิจผักสดต้องเพิ่มแกน operation: ล็อตของสด, รอบส่ง, แพ็กสินค้า, ปรับน้ำหนักจริง, subscription และการเคลมความสด", WHITE, RGBColor(207, 226, 207), DARK_GREEN, "🌿", 20, 14)
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_header(slide, "1) ภาพรวมสถานะปัจจุบัน", "โปรเจคมี foundation ด้าน LINE / CRM / Mini App แล้ว แต่ยังไม่ครบสำหรับธุรกิจผักสด")
    add_leaf_decor(slide)
    bullet_box(slide, 0.65, 1.45, 3.65, "มีแล้ว / ใช้ต่อได้", [
        "LINE OA + Inbox v2 สำหรับแชทลูกค้า",
        "Mini App / Shop สำหรับหน้าสั่งสินค้า",
        "Customer profile, tag, note, chat status",
        "Order management พื้นฐาน",
        "Admin menu/settings/shop content",
        "Notification / broadcast บางส่วน"
    ], GREEN, SOFT_GREEN, "✓")
    bullet_box(slide, 4.85, 1.45, 3.65, "ต้องปรับให้ตรงผักสด", [
        "สินค้าไม่ใช่ SKU ธรรมดา ต้องมีหน่วย/น้ำหนัก/เกรด",
        "Order ต้องรองรับน้ำหนักจริงหลังชั่ง",
        "Mini App ต้องเน้นรอบส่งและความสด",
        "Dashboard ต้องเห็นคิวแพ็กและของสดคงเหลือ",
        "แจ้งเตือนต้องผูกกับสถานะจัดส่ง"
    ], ORANGE, SOFT_ORANGE, "↺")
    bullet_box(slide, 9.05, 1.45, 3.65, "ยังขาดหลัก", [
        "Subscription Engine กล่องผัก",
        "Pre-order / Cut-off time",
        "Fresh stock lot / shelf life",
        "Packing dashboard",
        "Delivery zone / driver / tracking",
        "Proof of delivery + freshness claim"
    ], RED, SOFT_RED, "!")
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_header(slide, "2) โมเดลเป้าหมายสำหรับธุรกิจผักสด", "ลูกค้าสั่งผ่าน LINE/Mini App แล้วไหลเข้า order → fresh stock → packing → delivery → tracking")
    add_leaf_decor(slide)
    add_textbox(slide, 0.28, 1.55, 0.36, 0.52, "ลูกค้า", 11, GREEN, True)
    add_textbox(slide, 0.22, 3.05, 0.50, 0.70, "ระบบร้าน", 11, BLUE, True)
    add_textbox(slide, 0.18, 5.05, 0.58, 0.70, "จัดส่ง", 11, RED, True)
    cards = [
        (1.0, 1.35, 2.1, 0.7, "LINE App", "Rich Menu / Chat", SOFT_GREEN, GREEN, GREEN, "LINE"),
        (3.55, 1.35, 2.1, 0.7, "สั่งครั้งเดียว", "One-time order", MINT, GREEN, GREEN, "🛒"),
        (6.1, 1.35, 2.1, 0.7, "Subscription", "กล่องผักรายสัปดาห์", MINT, GREEN, GREEN, "🎁"),
        (8.65, 1.35, 2.1, 0.7, "Pre-order", "เลือกวันส่งล่วงหน้า", SOFT_ORANGE, ORANGE, ORANGE, "📅"),
        (1.0, 2.85, 2.1, 0.75, "Mini App", "เลือกผัก / ตะกร้า", SOFT_PURPLE, PURPLE, PURPLE, "📱"),
        (3.55, 2.85, 2.1, 0.75, "Order Mgmt", "รับ / ยืนยัน / แก้ไข", SOFT_BLUE, BLUE, BLUE, "📋"),
        (6.1, 2.85, 2.1, 0.75, "Fresh Stock", "ล็อต / เกรด / น้ำหนัก", SOFT_BLUE, BLUE, BLUE, "🌿"),
        (8.65, 2.85, 2.1, 0.75, "Payment", "PromptPay / COD", SOFT_ORANGE, ORANGE, ORANGE, "💳"),
        (1.0, 4.7, 2.1, 0.75, "Packing", "คัด / ชั่ง / แพ็ก", SOFT_RED, RED, RED, "📦"),
        (3.55, 4.7, 2.1, 0.75, "Delivery Slot", "รอบส่ง / โซน", SOFT_RED, RED, RED, "🚚"),
        (6.1, 4.7, 2.1, 0.75, "Driver", "รับงาน / อัปเดต", SOFT_RED, RED, RED, "🛵"),
        (8.65, 4.7, 2.1, 0.75, "Live Tracking", "ส่งลิงก์ใน LINE", SOFT_RED, RED, RED, "📍"),
        (11.0, 4.7, 1.7, 0.75, "POD", "รูป / ยืนยัน", SOFT_RED, RED, RED, "📷"),
    ]
    for c in cards:
        add_card(slide, *c, title_size=13, body_size=9.5)
    for x1, y1, x2, y2 in [(2.05, 2.05, 2.05, 2.85), (4.6, 2.05, 4.6, 2.85), (7.15, 2.05, 4.6, 2.85), (9.7, 2.05, 4.6, 2.85), (3.1, 3.22, 3.55, 3.22), (5.65, 3.22, 6.1, 3.22), (8.2, 3.22, 8.65, 3.22), (4.6, 3.6, 2.05, 4.7), (2.05, 5.45, 3.55, 5.1), (5.65, 5.1, 6.1, 5.1), (8.2, 5.1, 8.65, 5.1), (10.75, 5.1, 11.0, 5.1)]:
        add_arrow(slide, x1, y1, x2, y2, GRAY, 1.4)
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_header(slide, "3) Gap Matrix เทียบกับโมเดลในภาพ", "สถานะระบบตามกล่องสำคัญใน business model")
    headers = ["ระบบ", "สถานะ", "สิ่งที่ต้องทำเพิ่ม"]
    x_cols = [0.65, 4.35, 6.0]
    widths = [3.45, 1.35, 6.85]
    for x, w, htxt in zip(x_cols, widths, headers):
        add_card(slide, x, 1.35, w, 0.42, htxt, fill=DARK_GREEN, line=DARK_GREEN, title_color=WHITE, radius=False, title_size=12)
    rows = [
        ("LINE OA / Rich Menu", "มีแล้ว", "ปรับเมนูให้ตรงผักสด: สั่งผัก, กล่องผัก, ติดตาม, เคลม"),
        ("Mini App / Shop", "มีบางส่วน", "เพิ่มหน่วยขาย, น้ำหนัก, เกรด, รอบส่ง, cutoff time"),
        ("Customer Profile", "มีแล้ว", "เพิ่ม food preference: แพ้ผัก/ไม่กิน/ที่อยู่ประจำ"),
        ("Order Management", "มีพื้นฐาน", "เพิ่ม order state สำหรับคัด/ชั่ง/แพ็ก/พร้อมส่ง"),
        ("Payment", "มีบางส่วน", "รองรับปรับยอดหลังชั่งจริง, COD/เครดิต/แพ็กเกจ"),
        ("Subscription Engine", "ขาด", "กล่องผักรายสัปดาห์/รายเดือน, pause, renew, skip"),
        ("Delivery API / Zone", "ขาด", "พื้นที่ส่ง, รอบส่ง, capacity, ค่าส่ง, route batch"),
        ("Driver App / Tracking / POD", "ขาด", "รับงาน, live status, รูปยืนยัน, ส่งสำเร็จ"),
    ]
    y = 1.88
    for i, (system, status, action) in enumerate(rows):
        fill = RGBColor(255, 255, 255) if i % 2 == 0 else RGBColor(249, 252, 248)
        add_card(slide, x_cols[0], y, widths[0], 0.46, system, fill=fill, line=RGBColor(225, 232, 225), title_color=BLACK, radius=False, title_size=10.5)
        if status == "มีแล้ว":
            sf, sl, sc = SOFT_GREEN, GREEN, DARK_GREEN
        elif status == "มีบางส่วน":
            sf, sl, sc = SOFT_ORANGE, ORANGE, ORANGE
        else:
            sf, sl, sc = SOFT_RED, RED, RED
        add_card(slide, x_cols[1], y, widths[1], 0.46, status, fill=sf, line=sl, title_color=sc, radius=False, title_size=10)
        add_card(slide, x_cols[2], y, widths[2], 0.46, action, fill=fill, line=RGBColor(225, 232, 225), title_color=BLACK, radius=False, title_size=9.5)
        y += 0.55
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_header(slide, "4) ระบบเฉพาะผักสดที่ยังไม่มี", "ส่วนนี้ต่างจาก e-commerce ทั่วไป เพราะสินค้าเน่าเสียง่ายและน้ำหนักไม่เป๊ะ")
    modules = [
        ("Fresh Product Catalog", "หน่วยขาย: กก. / กรัม / แพ็ก / กล่อง\nเกรด: A, Organic, Hydroponic\nรูปจริง + ขนาดแพ็ก", SOFT_GREEN, GREEN, "🌿"),
        ("Fresh Stock Lot", "ล็อตรับเข้า / วันเก็บเกี่ยว\nShelf life / สินค้าใกล้หมด\nกัน oversell ของสด", SOFT_GREEN, GREEN, "🧺"),
        ("Weight Adjustment", "สั่ง 1 กก. แต่ชั่งจริง 0.95/1.05\nปรับยอดก่อนจ่ายหรือทำเครดิต\nบันทึกน้ำหนักจริงต่อรายการ", SOFT_ORANGE, ORANGE, "⚖"),
        ("Packing Dashboard", "คิวแพ็ก / คัดคุณภาพ\nชั่ง / แพ็ก / พร้อมส่ง\nพิมพ์ใบจัดของ", SOFT_BLUE, BLUE, "📦"),
        ("Freshness Claim", "ลูกค้าแนบรูปผักช้ำ/เสีย\nคืนเครดิต / ส่งใหม่ / คืนเงิน\nผูกกับ order และ lot", SOFT_RED, RED, "📷"),
        ("Waste Report", "ของเหลือ / ของเสีย / markdown\nกำไรจริงต่อ SKU\nช่วยวางแผนซื้อของเข้ารอบถัดไป", SOFT_PURPLE, PURPLE, "📊"),
    ]
    pos = [(0.7, 1.38), (4.7, 1.38), (8.7, 1.38), (0.7, 3.95), (4.7, 3.95), (8.7, 3.95)]
    for (title, body, fill, line, icon), (x, y) in zip(modules, pos):
        add_card(slide, x, y, 3.55, 1.95, title, body, fill=fill, line=line, title_color=line, icon=icon, title_size=15, body_size=10.5)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_header(slide, "5) Order Flow ใหม่สำหรับผักสด", "เพิ่มจุดคัด/ชั่ง/แพ็ก ก่อนส่ง เพื่อรองรับธรรมชาติของของสด")
    flow = [
        ("รับออเดอร์", "LINE / Mini App", SOFT_GREEN, GREEN),
        ("ตรวจสต็อก", "lot / shelf life", SOFT_BLUE, BLUE),
        ("จองของ", "reserve quantity", SOFT_BLUE, BLUE),
        ("คัดผัก", "quality check", SOFT_ORANGE, ORANGE),
        ("ชั่งจริง", "actual weight", SOFT_ORANGE, ORANGE),
        ("ปรับยอด", "เพิ่ม/ลด/เครดิต", SOFT_PURPLE, PURPLE),
        ("แพ็ก", "label / bag / box", SOFT_RED, RED),
        ("พร้อมส่ง", "assign rider", SOFT_RED, RED),
    ]
    x = 0.65
    y = 2.1
    for i, (title, body, fill, line) in enumerate(flow):
        add_card(slide, x + i * 1.55, y, 1.25, 0.95, title, body, fill=fill, line=line, title_color=line, title_size=11, body_size=8.3)
        if i < len(flow) - 1:
            add_arrow(slide, x + i * 1.55 + 1.25, y + 0.47, x + (i + 1) * 1.55, y + 0.47, GRAY, 1.1)
    add_card(slide, 0.75, 4.0, 3.75, 1.1, "Order states ที่ควรเพิ่ม", "new → confirmed → reserved → picking → weighed → payment_adjusted → packed → out_for_delivery → delivered / claimed", WHITE, GREEN, DARK_GREEN, "📋", 14, 10)
    add_card(slide, 4.9, 4.0, 3.75, 1.1, "ข้อมูลที่ต้องบันทึก", "ordered_qty, reserved_qty, actual_weight, unit_price, final_price, lot_id, packer_id, quality_note", WHITE, BLUE, BLUE, "🧾", 14, 10)
    add_card(slide, 9.05, 4.0, 3.25, 1.1, "เหตุผล", "ลดปัญหาของหมด, น้ำหนักไม่ตรง, ส่งผิด, เคลมยาก", WHITE, ORANGE, ORANGE, "!", 14, 10)
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_header(slide, "6) Subscription Box สำหรับกล่องผัก", "หัวใจของ recurring revenue ในธุรกิจผักสด")
    add_card(slide, 0.75, 1.45, 2.55, 1.1, "แพ็กเกจ", "กล่องเล็ก / กลาง / ใหญ่\nรายสัปดาห์ / รายเดือน", SOFT_GREEN, GREEN, DARK_GREEN, "🎁")
    add_card(slide, 3.75, 1.45, 2.55, 1.1, "Preference", "ไม่เอาผักบางชนิด\nแพ้อาหาร / organic only", SOFT_BLUE, BLUE, BLUE, "👤")
    add_card(slide, 6.75, 1.45, 2.55, 1.1, "รอบส่ง", "เลือกวัน / เวลา\nskip / pause / resume", SOFT_ORANGE, ORANGE, ORANGE, "📅")
    add_card(slide, 9.75, 1.45, 2.55, 1.1, "Billing", "ต่ออายุ / แจ้งเตือน\nชำระรายรอบ", SOFT_PURPLE, PURPLE, PURPLE, "💳")
    for x1, x2 in [(3.3, 3.75), (6.3, 6.75), (9.3, 9.75)]:
        add_arrow(slide, x1, 2.0, x2, 2.0, GRAY, 1.3)
    bullet_box(slide, 0.9, 3.25, 5.1, "ฟีเจอร์ขั้นต่ำ", [
        "สร้าง plan: weekly/monthly, price, included items/value",
        "subscription status: active, paused, cancelled, expired",
        "next_delivery_date + cut-off time",
        "เปลี่ยนกล่อง / skip รอบ / pause ชั่วคราว",
        "แจ้งเตือนก่อนตัดรอบและก่อนส่ง"
    ], GREEN, SOFT_GREEN, "✓")
    bullet_box(slide, 6.7, 3.25, 5.1, "ความเสี่ยงที่ต้องออกแบบ", [
        "ของสดหมดหรือราคาขึ้น: ต้องมี substitute rule",
        "ลูกค้าไม่อยู่บ้าน: reschedule / leave at door",
        "น้ำหนักจริงคลาดเคลื่อน: credit / adjust",
        "ผักเสีย: claim credit ผูกกับ lot"
    ], RED, SOFT_RED, "!")
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_header(slide, "7) Delivery Model ที่ยังขาด", "จากภาพต้นแบบ จุดแดงด้านล่างคือส่วนที่โปรเจคเรายังต้องสร้าง")
    steps = [
        ("Delivery Zone", "BKK / ปริมณฑล\nขั้นต่ำ / ค่าส่ง", 0.85),
        ("Delivery Slot", "รอบเช้า / บ่าย / เย็น\ncapacity ต่อรอบ", 3.1),
        ("Driver App", "รับงาน / map / โทรลูกค้า\nอัปเดตสถานะ", 5.35),
        ("Live Tracking", "ลิงก์ใน LINE\nETA / สถานะ", 7.6),
        ("Proof of Delivery", "ถ่ายรูป / ลายเซ็น\nส่งสำเร็จ", 9.85),
    ]
    for title, body, x in steps:
        add_card(slide, x, 2.0, 1.85, 1.15, title, body, SOFT_RED, RED, RED, "🚚", 12, 8.8)
    for x in [2.7, 4.95, 7.2, 9.45]:
        add_arrow(slide, x, 2.57, x + 0.38, 2.57, RED, 1.4)
    add_card(slide, 0.9, 4.3, 3.4, 1.15, "ข้อมูล delivery ที่ต้องมี", "zone_id, slot_id, route_batch_id, driver_id, status, ETA, proof_photo_url", WHITE, RED, RED, "🧾", 14, 10)
    add_card(slide, 4.85, 4.3, 3.4, 1.15, "LINE notification", "ยืนยันออเดอร์ → พร้อมส่ง → กำลังจัดส่ง → ส่งสำเร็จ → ขอรีวิว/เคลม", WHITE, GREEN, GREEN, "LINE", 14, 10)
    add_card(slide, 8.8, 4.3, 3.2, 1.15, "เริ่มแบบง่ายก่อน", "ยังไม่ต้องทำ GPS จริงทันที เริ่มจาก status tracking + proof photo ได้", WHITE, ORANGE, ORANGE, "MVP", 14, 10)
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_header(slide, "8) Data Model ที่ควรเพิ่ม", "ตารางหลักที่จำเป็นสำหรับผักสดและ subscription")
    left_tables = [
        ("fresh_product_units", "unit, min_qty, step_qty"),
        ("fresh_product_grades", "grade, claim_rule"),
        ("fresh_stock_lots", "product_id, received_at, harvest_at, expires_at"),
        ("stock_movements", "in/out/waste/reserve"),
        ("order_item_weights", "ordered_qty, actual_weight, final_price"),
    ]
    right_tables = [
        ("subscription_plans", "box size, frequency, price"),
        ("customer_subscriptions", "status, next_delivery, preference"),
        ("delivery_zones", "area, fee, min_order"),
        ("delivery_slots", "date, time, capacity"),
        ("delivery_tasks", "driver, route, status, proof"),
        ("freshness_claims", "order_id, photo, resolution"),
    ]
    add_card(slide, 0.8, 1.4, 5.55, 0.5, "Fresh Inventory / Order", fill=SOFT_GREEN, line=GREEN, title_color=DARK_GREEN, title_size=14)
    y = 2.0
    for name, desc in left_tables:
        add_card(slide, 0.9, y, 5.35, 0.43, name + "  —  " + desc, fill=WHITE, line=RGBColor(220, 232, 220), title_color=BLACK, radius=False, title_size=9.5)
        y += 0.53
    add_card(slide, 7.0, 1.4, 5.55, 0.5, "Subscription / Delivery / Claim", fill=SOFT_ORANGE, line=ORANGE, title_color=ORANGE, title_size=14)
    y = 2.0
    for name, desc in right_tables:
        add_card(slide, 7.1, y, 5.35, 0.43, name + "  —  " + desc, fill=WHITE, line=RGBColor(238, 226, 210), title_color=BLACK, radius=False, title_size=9.5)
        y += 0.53
    add_card(slide, 0.95, 5.6, 11.45, 0.85, "หมายเหตุ", "ควรเริ่มจากตารางที่รองรับ MVP: product units, stock lots, order item weights, delivery slots, delivery tasks แล้วค่อยเพิ่ม subscription/claim เชิงลึก", SOFT_BLUE, BLUE, BLUE, "⚙", 14, 10.5)
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_header(slide, "9) Roadmap แนะนำ", "ทำทีละเฟสเพื่อลดความเสี่ยง และใช้ของเดิมให้มากที่สุด")
    phases = [
        ("Phase 1", "Fresh Shop MVP", "สินค้า/หน่วยขาย/เกรด\nรอบส่งพื้นฐาน\nMini App copy + menu", GREEN),
        ("Phase 2", "Fresh Operation", "Fresh stock lot\nชั่งน้ำหนักจริง\nPacking dashboard", BLUE),
        ("Phase 3", "Subscription", "กล่องผัก\nรอบ recurring\nแจ้งเตือนต่ออายุ/skip", ORANGE),
        ("Phase 4", "Delivery", "Driver task\nTracking link\nProof of delivery", RED),
        ("Phase 5", "Optimization", "Claim/waste report\nForecast demand\nMargin analytics", PURPLE),
    ]
    x = 0.65
    for i, (p, title, body, color) in enumerate(phases):
        add_card(slide, x + i * 2.5, 1.55, 2.15, 2.0, p, title + "\n\n" + body, RGBColor(255, 255, 255), color, color, "●", 13, 9.2)
        if i < len(phases) - 1:
            add_arrow(slide, x + i * 2.5 + 2.15, 2.55, x + (i + 1) * 2.5, 2.55, GRAY, 1.2)
    add_card(slide, 0.9, 4.4, 3.55, 1.25, "MVP แรกที่ควรทำ", "Fresh Product Catalog + Delivery Slot + Packing status", SOFT_GREEN, GREEN, DARK_GREEN, "1", 15, 12)
    add_card(slide, 4.9, 4.4, 3.55, 1.25, "ห้ามข้าม", "Stock lot + actual weight เพราะเป็นจุดต่างของผักสด", SOFT_ORANGE, ORANGE, ORANGE, "!", 15, 12)
    add_card(slide, 8.9, 4.4, 3.2, 1.25, "ทำหลัง MVP", "GPS tracking จริง / route optimization / demand forecast", SOFT_BLUE, BLUE, BLUE, "→", 15, 12)
    add_footer(slide, 10)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    add_header(slide, "10) Priority Backlog", "รายการงานที่ควรสร้างก่อน-หลังตามผลกระทบธุรกิจ")
    headers = ["Priority", "Module", "Outcome", "Effort"]
    x_cols = [0.75, 2.1, 5.15, 10.75]
    widths = [1.05, 2.75, 5.25, 1.55]
    for x, w, htxt in zip(x_cols, widths, headers):
        add_card(slide, x, 1.35, w, 0.42, htxt, fill=DARK_GREEN, line=DARK_GREEN, title_color=WHITE, radius=False, title_size=11)
    rows = [
        ("P0", "Fresh unit + grade", "เริ่มขายผักสดแบบถูกหน่วยและสื่อสารความสดได้", "M"),
        ("P0", "Delivery slot / cutoff", "รับออเดอร์ตามรอบส่งจริง ลดปัญหาส่งไม่ทัน", "M"),
        ("P1", "Stock lot + expiry", "รู้ของสดเหลือ/ใกล้เสีย/กันขายเกิน", "L"),
        ("P1", "Packing dashboard", "ทีมแพ็กเห็นคิวและสถานะได้ชัด", "M"),
        ("P1", "Actual weight adjustment", "แก้ปัญหาน้ำหนักไม่เป๊ะและยอดชำระ", "M"),
        ("P2", "Subscription box", "สร้างรายได้ประจำและกล่องผักรายสัปดาห์", "L"),
        ("P2", "Driver + POD", "ปิดงานส่งและลดข้อโต้แย้ง", "M"),
        ("P3", "Claim + waste report", "ควบคุมคุณภาพและ margin ของสด", "M"),
    ]
    y = 1.88
    for i, row in enumerate(rows):
        priority, module, outcome, effort = row
        fill = RGBColor(255, 255, 255) if i % 2 == 0 else RGBColor(249, 252, 248)
        pcol = RED if priority == "P0" else ORANGE if priority == "P1" else BLUE if priority == "P2" else GRAY
        add_card(slide, x_cols[0], y, widths[0], 0.44, priority, fill=fill, line=RGBColor(225, 232, 225), title_color=pcol, radius=False, title_size=10.5)
        add_card(slide, x_cols[1], y, widths[1], 0.44, module, fill=fill, line=RGBColor(225, 232, 225), title_color=BLACK, radius=False, title_size=9.5)
        add_card(slide, x_cols[2], y, widths[2], 0.44, outcome, fill=fill, line=RGBColor(225, 232, 225), title_color=BLACK, radius=False, title_size=9.5)
        add_card(slide, x_cols[3], y, widths[3], 0.44, effort, fill=fill, line=RGBColor(225, 232, 225), title_color=GRAY, radius=False, title_size=10)
        y += 0.52
    add_footer(slide, 11)

    # Slide 12
    slide = prs.slides.add_slide(blank)
    add_header(slide, "11) สรุปสำหรับการตัดสินใจ", "ถ้าจะไปโมเดลผักสด ต้องเพิ่ม operation layer ไม่ใช่แค่เปลี่ยนหน้าร้าน")
    add_card(slide, 0.85, 1.35, 3.4, 1.35, "สิ่งที่มีแล้ว", "LINE OA, Inbox, Mini App, Customer, Order, Admin\nใช้เป็นฐานได้ทันที", SOFT_GREEN, GREEN, DARK_GREEN, "✓", 17, 12)
    add_card(slide, 4.95, 1.35, 3.4, 1.35, "สิ่งที่ต้องสร้าง", "Fresh stock, packing, delivery slot, subscription, tracking, claim", SOFT_ORANGE, ORANGE, ORANGE, "!", 17, 12)
    add_card(slide, 9.05, 1.35, 3.2, 1.35, "คำแนะนำ", "ทำ MVP ด้วย Fresh Catalog + Slot + Packing ก่อน แล้วค่อยต่อ subscription/driver", SOFT_BLUE, BLUE, BLUE, "→", 17, 12)
    add_textbox(slide, 1.05, 3.55, 11.2, 0.65, "Core principle: จาก e-commerce ทั่วไป → Fresh operation platform", 25, DARK_GREEN, True, PP_ALIGN.CENTER)
    add_textbox(slide, 1.7, 4.42, 10.0, 0.48, "ผักสดต้องควบคุม 4 เรื่อง: ความสด, น้ำหนักจริง, รอบส่ง, เคลมคุณภาพ", 18, ORANGE, True, PP_ALIGN.CENTER)
    add_footer(slide, 12)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    create_deck()
